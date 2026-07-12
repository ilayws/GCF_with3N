#!/usr/bin/env python3
"""Generate a plain black-on-white summary deck of the GCF generators + GENIE FSI."""
from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()  # default 10 x 7.5 in, default template
BLANK = prs.slide_layouts[6]

def slide(title, bullets):
    s = prs.slides.add_slide(BLANK)
    # Title
    tb = s.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(26); r.font.bold = True
    # Body
    bb = s.shapes.add_textbox(Inches(0.4), Inches(1.15), Inches(9.2), Inches(6.0))
    bf = bb.text_frame
    bf.word_wrap = True
    first = True
    for text, lvl in bullets:
        p = bf.paragraphs[0] if first else bf.add_paragraph()
        first = False
        p.level = lvl
        prefix = "" if lvl == 0 else ("  " * lvl)
        bullet = "• " if lvl == 0 else "- "
        run = p.add_run(); run.text = prefix + bullet + text
        run.font.size = Pt(16) if lvl == 0 else Pt(13)
        p.space_after = Pt(4)
    return s

def title_slide(title, subtitle_lines):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(8.8), Inches(2))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(32); r.font.bold = True
    for line in subtitle_lines:
        p = tf.add_paragraph(); r = p.add_run(); r.text = line
        r.font.size = Pt(16)
    return s

# ---------------------------------------------------------------- title
title_slide(
    "GCF Quasi-Elastic Event Generators + GENIE FSI",
    ["Mean-field (1N), 2N SRC, 3N SRC generators and the GENIE final-state-interaction layer",
     "Technical functional summary (genQE_FSI, genQE_3N)"])

# ---------------------------------------------------------------- common engine
slide("Common Engine: Weighted Monte Carlo", [
    ("All four generators share one importance-sampling engine.", 0),
    ("Each trial starts weight=1; every sampled variable multiplies weight by its interval volume.", 0),
    ("Physics factors multiplied in: structure/density function, e-N cross section, energy-conservation Jacobians.", 0),
    ("Kinematically invalid trial -> weight=0, discarded. Surviving events + weights -> ROOT TTree (weight-summed cross sections).", 0),
    ("Beam: electron along z, fixed energy.", 0),
    ("e-N cross section sigma_eN: De Forest CC1 (default) or CC2 off-shell, or on-shell Rosenbluth.", 0),
    ("Form factors: Kelly (default); dipole or Ye (Galster for neutron electric) optional.", 0),
    ("Optional: electron radiative corrections (peaking approx.) and Coulomb correction. Off by default.", 0),
])

# ---------------------------------------------------------------- mean field
slide("1. Mean-Field (single-nucleon) generator", [
    ("Quasi-elastic knockout of one uncorrelated nucleon from a Fermi gas. (generate_event_MF)", 0),
    ("Pick struck-nucleon isospin p/n 50/50, weight x2; no correlated partner.", 0),
    ("Sample initial momentum from a Fermi gas: isotropic direction, |p| uniform in [0, kF].", 0),
    ("Global FG: uniform-sphere n(p)=3/(4 pi kF^3).", 1),
    ("Local FG: Thomas-Fermi kF(r)=(3 pi^2 rho(r)/2)^(1/3) from GENIE density, tabulated n(p).", 1),
    ("Weight folds in phase-space volume x p^2 n(p)/(2 pi)^3.", 1),
    ("(A-1) residual: tabulated ground-state mass, recoils with -p1.", 0),
    ("Struck nucleon off-shell: E1 = mA - E(A-1) by energy conservation.", 0),
    ("Lepton vertex solved (shared solver); FSI transports the nucleon through (A-1).", 0),
])

# ---------------------------------------------------------------- 2N
slide("2. Two-Nucleon (2N) SRC generator  (GCF)", [
    ("Scattering off one nucleon of a short-range-correlated pair. (generate_event)", 0),
    ("Assign pair isospin lead+recoil 50/50, weight x4 (pp/pn/np/nn).", 0),
    ("(A-2) residual: tabulated mass, optional excitation E* (fixed or Gaussian).", 0),
    ("Decay function: CM momentum from 3-D Gaussian (sigma_CM ~ 0.10-0.15 GeV/c).", 0),
    ("Relative momentum: direction uniform, magnitude uniform in [pRel_min, pRel_max].", 0),
    ("Weight x phase space x pRel^2 x S(pRel)/(2 pi)^3, with S = C_ij |phi_ij(pRel)|^2", 1),
    ("(contact coefficient C x universal pair function phi; tabulated per NN model: AV18, AV4', N2LO, ...).", 1),
    ("Momenta: p1=CM/2+pRel (struck), pRec=CM/2-pRel; (A-2) recoils -(p1+pRec). Struck off-shell: E1=mA-E(A-2)-Erec.", 0),
])

slide("2. 2N generator - lepton vertex & variants", [
    ("Electron vertex (core solve): sample Q^2 in [Q2min, min(Q2max, kinematic limit)] and azimuth phi_k.", 0),
    ("Solve quadratic for outgoing electron transverse momentum enforcing energy-momentum conservation.", 0),
    ("Up to two roots (extra x2 weight + random choice). Build q and p_lead = p1 + q.", 1),
    ("Divide weight by the energy-conservation delta-function Jacobian.", 0),
    ("Multiply by off-shell sigma_eN; apply optional radiation/Coulomb.", 0),
    ("Light-cone variant (generate_event_lightcone): samples pair in light-cone fractions alpha, solves for omega directly.", 0),
])

# ---------------------------------------------------------------- 3N
slide("3. Three-Nucleon (3N) SRC generator", [
    ("e + A -> e' + N1 + N2 + N3 + (A-3)*; scattering off one nucleon of a correlated triplet. (QEGenerator_3N)", 0),
    ("Key change vs 2N: contact x universal function -> pre-computed 3-body density matrix (ab-initio).", 0),
    ("ppn triplets only: start all protons, flip one (uniform among 3) to neutron, weight x3.", 0),
    ("Target/residual masses configurable (code default 12C -> 9Be; PHYSICS note example 4He -> p). Residual may carry E*.", 0),
    ("3-body CM momentum from 3-D Gaussian (code default sigma_CM=0.15 GeV/c; note says 0.055). Each nucleon +P_CM/3; residual -P_CM.", 0),
    ("Orientation: 3 Euler angles uniform, weight x 8 pi^2.", 0),
    ("Internal Jacobi vars sampled in 'primed' wavefunction frame: p_a=|neutron|, p_b=|pp rel /2|, theta_ab; weight x volume.", 0),
])

slide("3. 3N generator - momenta, vertex, weight", [
    ("Build v1=p_a, v2=-p_a/2+p_b, v3=-p_a/2-p_b; Euler-rotate; add P_CM/3. Inverse permutation sets struck=N1.", 0),
    ("Reject any |p| outside [0.25, 5] GeV/c.", 1),
    ("Energies: residual + 2 spectators on-shell; struck E1 = mA - E(A-3) - E2 - E3 (off-shell).", 0),
    ("Electron vertex: sample direction (cos theta_k in [cos40,cos5], phi_k); solve E_k in CLOSED FORM (no quadratic).", 0),
    ("Fold in energy-conservation delta-function Jacobian.", 0),
    ("Weight x sigma_eN (CC1/Kelly) x rho_3N(theta_ab, k_cm, k_rel), trilinear interp of 61x101x101 table.", 0),
    ("NN-interaction tables: AV8', chiral N2LO, G3/N3LO, AV4'. Fixed norm constants C=0.12, t=2, 1/2 (2 pi)^-7.", 0),
    ("Driver auto-retries until target # of weight>0 events (low 9-D acceptance).", 0),
])

# ---------------------------------------------------------------- FSI 1
slide("4. GENIE FSI layer - setup & geometry", [
    ("Post-processing of the bare PWIA event. Single compile flag USE_FSI; backend = local GENIE R-3.06.02.", 0),
    ("Residual medium traversed: A-1 (MF), A-2 (2N), A-3 (3N).", 0),
    ("Skip FSI if A_res < 3 (cascade not validated, NaN splines); one-time warning for A_res < 7.", 0),
    ("Common vertex for the correlated nucleons, rejection-sampled from rho^n(r) r^2:", 0),
    ("n=1 (MF single), n=2 (pair), n=3 (triplet); rho from GENIE density model.", 1),
])

slide("4. GENIE FSI layer - cascade & cuts", [
    ("Per nucleon: build GENIE record in lepton-nucleus mode (dummy high-E electron probe -> GENIE trusts our vertex).", 0),
    ("Target = residual ion (A_res,Z_res); nucleon = 'hadron in nucleus' with its 4-momentum + sampled position.", 0),
    ("Run INTRANUKE 2018 cascade: hN (full cascade, default) or hA (effective single step). GENIE RNG reseeded from project RNG.", 0),
    ("Follow chain to surviving stable nucleon (highest-E if several): updates 4-momentum and species (captures charge exchange).", 0),
    ("Other stable products recorded as secondaries (pi+/pi0/pi- counted). Full absorption -> weight=0.", 0),
    ("Pauli blocking: any final |p| < Fermi momentum (default 220 MeV/c) -> weight=0.", 0),
    ("Pre-FSI momenta stored alongside post-FSI; FSI can be disabled for pure PWIA.", 0),
])

# ---------------------------------------------------------------- caveats
slide("Orchestration & Caveats", [
    ("Each driver loops over trials, runs generator + (optional) GENIE FSI, writes weight>0 events to a ROOT TTree.", 0),
    ("3N driver loops until a target accepted count; 2N/MF drivers loop a fixed number of trials.", 0),
    ("Caveat: 3N PHYSICS.tex is partly outdated vs code - sigma_CM (0.15 vs 0.055), default target (12C vs 4He), removed coord Jacobian.", 0),
    ("Caveat: 3N constants C, t are fixed numerical factors, not first-principles -> absolute 3N rates are normalization-dependent; shapes are not.", 0),
])

prs.save("GCF_generators_summary.pptx")
print("wrote GCF_generators_summary.pptx with", len(prs.slides._sldIdLst), "slides")
